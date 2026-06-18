# -*- coding: utf-8 -*-
"""Genera presentacion-infra-seguridad.pptx con la identidad visual de Trackt.
Tema dark, brand violeta #6152e8, fuentes Arial/Consolas, colores de estado de tickets."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Paleta Trackt (dark) ----------
BG        = RGBColor(0x0B, 0x0B, 0x0D)   # neutral-950 ~
SURFACE   = RGBColor(0x15, 0x15, 0x18)   # card
SURFACE_2 = RGBColor(0x1D, 0x1D, 0x21)   # card alt
BORDER    = RGBColor(0x2A, 0x2A, 0x30)
FG        = RGBColor(0xF5, 0xF5, 0xF5)   # neutral-100
MUTED     = RGBColor(0xA0, 0xA0, 0xA8)   # neutral-400/500
FAINT     = RGBColor(0x70, 0x70, 0x78)

BRAND     = RGBColor(0x61, 0x52, 0xE8)   # brand-500
BRAND_LT  = RGBColor(0x9B, 0x8F, 0xFF)   # brand-300
BRAND_DK  = RGBColor(0x3D, 0x31, 0xB0)   # brand-700

BLUE   = RGBColor(0x60, 0xA5, 0xFA)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
GREEN  = RGBColor(0x4A, 0xDE, 0x80)
RED    = RGBColor(0xF8, 0x71, 0x71)
EMER   = RGBColor(0x34, 0xD3, 0x99)

F_SANS = "Arial"
F_MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def text(s, x, y, w, h, runs, size=18, color=FG, bold=False, font=F_SANS,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    if spacing: p.line_spacing = spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, st in runs:
        r = p.add_run(); r.text = t
        f = r.font
        f.size = Pt(st.get("size", size)); f.bold = st.get("bold", bold)
        f.italic = st.get("italic", italic)
        f.name = st.get("font", font)
        f.color.rgb = st.get("color", color)
    return tb

def bullets(s, x, y, w, h, items, size=16, gap=8, marker_color=BRAND):
    """items: list of (text, level) or (text, level, color)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        txt = it[0]; lvl = it[1]; col = it[2] if len(it) > 2 else FG
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.line_spacing = 1.05
        if lvl == 0:
            mk = p.add_run(); mk.text = "▸  "
            mk.font.name = F_SANS; mk.font.size = Pt(size); mk.font.bold = True
            mk.font.color.rgb = marker_color
            r = p.add_run(); r.text = txt
            r.font.name = F_SANS; r.font.size = Pt(size); r.font.color.rgb = col
        else:
            p.level = 1
            mk = p.add_run(); mk.text = "      –   "
            mk.font.name = F_SANS; mk.font.size = Pt(size-1); mk.font.color.rgb = FAINT
            r = p.add_run(); r.text = txt
            r.font.name = F_SANS; r.font.size = Pt(size-1); r.font.color.rgb = MUTED
    return tb

def chip(s, x, y, label, dot, w=Inches(1.95), h=Inches(0.42)):
    c = rect(s, x, y, w, h, fill=SURFACE_2, line=BORDER, line_w=1.0, rounded=True, radius=0.5)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.16), y + h/2 - Inches(0.07),
                           Inches(0.14), Inches(0.14))
    d.fill.solid(); d.fill.fore_color.rgb = dot; d.line.fill.background(); d.shadow.inherit = False
    text(s, x + Inches(0.4), y, w - Inches(0.45), h, label, size=11.5, color=FG,
         anchor=MSO_ANCHOR.MIDDLE, font=F_SANS)
    return c

def header(s, idx, kicker, title, kicker_color=BRAND_LT):
    rect(s, 0, 0, Inches(0.16), EMU_H, fill=BRAND)             # left brand spine
    text(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.3),
         kicker.upper(), size=12.5, color=kicker_color, bold=True)
    text(s, Inches(0.6), Inches(0.82), Inches(12), Inches(0.8),
         title, size=30, color=FG, bold=True, font=F_SANS)
    rect(s, Inches(0.62), Inches(1.62), Inches(0.9), Pt(3), fill=BRAND)
    footer(s, idx)

def footer(s, idx):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(7.06), Inches(0.13), Inches(0.13))
    d.fill.solid(); d.fill.fore_color.rgb = BRAND; d.line.fill.background(); d.shadow.inherit = False
    text(s, Inches(0.82), Inches(6.98), Inches(7), Inches(0.3),
         "Trackt · Infraestructura & Seguridad", size=10, color=FAINT)
    text(s, Inches(11.2), Inches(6.98), Inches(1.5), Inches(0.3),
         f"{idx:02d}", size=10, color=FAINT, align=PP_ALIGN.RIGHT)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def card(s, x, y, w, h, title, lines, accent=BRAND, title_color=FG):
    rect(s, x, y, w, h, fill=SURFACE, line=BORDER, line_w=1.0, rounded=True, radius=0.06)
    rect(s, x, y, w, Pt(4), fill=accent)
    text(s, x + Inches(0.22), y + Inches(0.18), w - Inches(0.4), Inches(0.4),
         title, size=15, color=title_color, bold=True)
    tb = s.shapes.add_textbox(x + Inches(0.22), y + Inches(0.62), w - Inches(0.4), h - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.line_spacing = 1.04
        r = p.add_run(); r.text = ln
        r.font.name = F_SANS; r.font.size = Pt(12.5); r.font.color.rgb = MUTED

# =====================================================================
# Slide 1 — Portada
# =====================================================================
s = slide()
rect(s, 0, 0, Inches(0.22), EMU_H, fill=BRAND)
# logo mark
lm = rect(s, Inches(0.9), Inches(1.15), Inches(0.62), Inches(0.62),
          fill=BRAND, rounded=True, radius=0.28)
text(s, Inches(0.9), Inches(1.15), Inches(0.62), Inches(0.62), "T",
     size=30, color=RGBColor(0xFF,0xFF,0xFF), bold=True, align=PP_ALIGN.CENTER,
     anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(1.7), Inches(1.18), Inches(6), Inches(0.6), "Trackt",
     size=40, color=FG, bold=True)
text(s, Inches(0.92), Inches(2.55), Inches(11), Inches(1.0),
     "Infraestructura & Seguridad", size=46, color=FG, bold=True)
rect(s, Inches(0.95), Inches(3.7), Inches(1.3), Pt(4), fill=BRAND)
text(s, Inches(0.95), Inches(3.95), Inches(11), Inches(0.5),
     "Plataforma de mantenimiento industrial — cómo está montado el sistema por dentro",
     size=16, color=MUTED)
# status chips row (identidad UI Trackt)
cx = Inches(0.95); cy = Inches(4.85)
for lab, dot in [("Pendiente", FAINT), ("Asignado", BLUE), ("En ejecución", AMBER),
                 ("Ejecutado", GREEN), ("Cerrado", EMER)]:
    chip(s, cx, cy, lab, dot, w=Inches(1.78))
    cx += Inches(1.9)
text(s, Inches(0.95), Inches(5.9), Inches(11), Inches(0.4),
     [("Equipo:  ", {"color": FAINT, "size": 14}),
      ("Rosio Ametller · Jaime Osorio · Ramón Hernández", {"color": FG, "size": 14, "bold": True})])
text(s, Inches(0.95), Inches(6.35), Inches(11), Inches(0.4),
     [("Duración objetivo:  ", {"color": FAINT, "size": 13}),
      ("~10 minutos", {"color": BRAND_LT, "size": 13, "bold": True})])
notes(s, "Presentación de ~10 min: primero infraestructura (tecnologías y dónde corre cada cosa) "
         "y luego seguridad (auth, roles, aislamiento, archivos). Trackt = sistema de gestión de "
         "taller / mantenimiento industrial: equipos, órdenes de trabajo, tickets, inventario y evidencias. (0:30)")

# =====================================================================
# Slide 2 — Agenda
# =====================================================================
s = slide()
header(s, 2, "Recorrido", "Agenda")
card(s, Inches(0.6), Inches(2.0), Inches(5.95), Inches(4.4),
     "01 · Infraestructura",
     ["Vista general — 3 capas desacopladas",
      "Frontend — Next.js 16 en Vercel",
      "Backend / API — NestJS 11 en Railway",
      "Base de datos — Supabase PostgreSQL",
      "",
      "Qué tecnología usamos y dónde corre cada pieza."],
     accent=BRAND)
card(s, Inches(6.78), Inches(2.0), Inches(5.95), Inches(4.4),
     "02 · Seguridad",
     ["Autenticación — JWT de Supabase",
      "Autorización — roles y mínimo privilegio",
      "Aislamiento multi-tenant",
      "Evidencias — signed URLs",
      "Secretos, red y RLS",
      "Proceso — Gitflow y PRs"],
     accent=BRAND_LT, title_color=BRAND_LT)
notes(s, "Mapa de la charla. Dos bloques: infraestructura y seguridad. La seguridad es donde "
         "pusimos más cuidado y le dedicamos más tiempo. (0:20)")

# =====================================================================
# Slide 3 — Divider Infraestructura
# =====================================================================
def divider(idx, num, title, sub):
    s = slide()
    rect(s, 0, 0, Inches(0.22), EMU_H, fill=BRAND)
    text(s, Inches(0.95), Inches(2.5), Inches(4), Inches(1.2), num,
         size=90, color=SURFACE_2, bold=True)
    text(s, Inches(0.98), Inches(3.75), Inches(11), Inches(0.9), title,
         size=44, color=FG, bold=True)
    rect(s, Inches(1.0), Inches(4.75), Inches(1.3), Pt(4), fill=BRAND)
    text(s, Inches(1.0), Inches(4.95), Inches(10.5), Inches(0.6), sub, size=16, color=MUTED)
    footer(s, idx)
    return s

s = divider(3, "01", "Infraestructura", "Tres capas independientes: cada una se despliega, escala y falla por separado.")
notes(s, "Pasamos al bloque de infraestructura. (transición, 0:05)")

# =====================================================================
# Slide 4 — Arquitectura general (3 cards + flujo)
# =====================================================================
s = slide()
header(s, 4, "Infraestructura", "Vista general de la arquitectura")
cy = Inches(2.1); cw = Inches(3.7); ch = Inches(3.05)
card(s, Inches(0.6), cy, cw, ch, "Frontend  ·  SSR",
     ["Next.js 16 (App Router)", "React 19", "TanStack Query", "shadcn/ui + Tailwind 4", "", "Deploy: Vercel"],
     accent=BLUE)
card(s, Inches(4.82), cy, cw, ch, "Backend  ·  API REST",
     ["NestJS 11 + TypeScript", "Prisma 6 (ORM)", "Auth JWT (Bearer)", "Validación global de DTOs", "", "Deploy: Railway"],
     accent=BRAND)
card(s, Inches(9.04), cy, cw, ch, "Datos  ·  Servicios",
     ["Supabase", "· PostgreSQL", "· Auth (JWT)", "· Storage (evidencias)", "", "Región: East US"],
     accent=EMER)
# arrows between
for ax in [Inches(4.42), Inches(8.64)]:
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, cy + ch/2 - Inches(0.16), Inches(0.34), Inches(0.32))
    a.fill.solid(); a.fill.fore_color.rgb = BRAND; a.line.fill.background(); a.shadow.inherit=False
bullets(s, Inches(0.62), Inches(5.55), Inches(12), Inches(1.2), [
    ("El frontend nunca habla directo con la base: todo pasa por la API.", 0),
    ("Comunicación HTTPS + token JWT en cada request → un único punto de entrada controlado.", 0),
], size=15)
notes(s, "Tres capas independientes. Front en Next.js (Vercel) sólo dibuja UI y protege rutas. "
         "API NestJS (Railway) tiene toda la lógica. Supabase da Postgres + Auth + Storage en uno. "
         "Lo desacoplado permite desplegar/escalar/fallar por separado. Primera decisión de seguridad: "
         "el front nunca toca la DB; todo pasa por la API por HTTPS con token. (1:15)")

# =====================================================================
# Slide 5 — Frontend
# =====================================================================
s = slide()
header(s, 5, "Infraestructura", "Frontend — Next.js 16 + Vercel")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("Next.js 16 (App Router) + React 19, rendering en servidor.", 0),
    ("Protección de rutas en el servidor:", 0),
    ("middleware valida la sesión (supabase.auth.getUser) en cada request", 1),
    ("redirige a /login si no hay sesión, antes de renderizar", 1),
    ("Sesión en cookies vía @supabase/ssr (no localStorage → menos XSS).", 0),
    ("TanStack Query para estado de servidor y caché.", 0),
    ("shadcn + Tailwind 4 · formularios validados con zod.", 0),
    ("authFetch adjunta el token como Bearer (auto-refresh + retry en 401).", 0),
    ("El build de producción falla si faltan las env vars de Supabase.", 0),
], size=15, gap=7)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Por qué importa",
     ["Sin lógica de negocio ni llaves secretas en el cliente.",
      "Protección server-side: no depende de 'esconder un botón'.",
      "Cookies > localStorage para la sesión.",
      "Token refrescado de forma transparente.",
      "",
      "producto/tract-front  (paquete trackt-front)"],
     accent=BLUE)
notes(s, "Next.js 16 con App Router en Vercel. A diferencia de una SPA, hay protección de rutas "
         "server-side: un middleware valida la sesión contra Supabase en cada request y manda a "
         "login si no estás autenticado. La sesión va en cookies (SSR), no en localStorage → menos "
         "XSS. Datos con TanStack Query, UI shadcn/Tailwind, forms con zod. authFetch pone el token "
         "Bearer y refresca/reintenta en 401. Si faltan env vars, el build de prod falla a propósito. (1:00)")

# =====================================================================
# Slide 6 — Backend / API
# =====================================================================
s = slide()
header(s, 6, "Infraestructura", "Backend / API — NestJS 11")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("NestJS 11 + TypeScript: arquitectura modular (un módulo por dominio).", 0),
    ("Prisma 6 como ORM → PostgreSQL (queries tipadas, anti-inyección).", 0),
    ("Validación global con ValidationPipe:", 0),
    ("whitelist → descarta campos no declarados en el DTO", 1),
    ("transform → convierte y valida tipos antes de la lógica", 1),
    ("Ningún dato entra a la API sin pasar por un filtro.", 0),
    ("Deploy en Railway (Node).", 0),
], size=15, gap=8)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Módulos de dominio",
     ["auth · tenant", "equipos · usuarios", "ordenes · tickets", "evidencias",
      "inventario", "notificaciones", "", "9 módulos aislados → reglas de seguridad por módulo"],
     accent=BRAND)
notes(s, "NestJS obliga a estructura modular: cada dominio (tickets, órdenes, inventario, "
         "evidencias…) es un módulo. Prisma como ORM con queries tipadas nos protege de inyección "
         "SQL. Pieza clave: ValidationPipe global — whitelist descarta campos no declarados, "
         "transform valida tipos antes de tocar la lógica. Nada entra sin filtro. Deploy en Railway. (1:15)")

# =====================================================================
# Slide 7 — Base de datos
# =====================================================================
s = slide()
header(s, 7, "Infraestructura", "Base de datos — Supabase PostgreSQL")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("PostgreSQL administrada por Supabase (región East US).", 0),
    ("Conexión vía pooler (modo transacción, 6543) en runtime.", 0),
    ("Conexión directa (5432) sólo para migraciones.", 0),
    ("sslmode=require → tráfico a la DB siempre cifrado.", 0),
    ("tenant_id en TODAS las tablas de negocio → multi-tenant desde el diseño.", 0, BRAND_LT),
], size=15, gap=9)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Modelo de datos",
     ["Tenant · Equipo", "OrdenTrabajo · Ticket", "Evidencia", "Repuesto · InventarioStock",
      "ReservaRepuesto", "MovimientoInventario", "Notificacion", "", "Cubre todo el flujo del taller"],
     accent=EMER)
notes(s, "Postgres gestionada por Supabase. En runtime usamos un pooler en modo transacción que "
         "aguanta mucha concurrencia; la conexión directa sólo para migraciones. Todo el tráfico "
         "a la DB va cifrado (SSL obligatorio). El modelo cubre el flujo completo del taller. "
         "Detalle clave: TODAS las tablas llevan tenant_id → multi-tenant por diseño. (1:00)")

# =====================================================================
# Slide 8 — Divider Seguridad
# =====================================================================
s = divider(8, "02", "Seguridad", "Defensa en profundidad: varias capas, no una sola.")
notes(s, "Entramos a seguridad, el núcleo de la charla. (transición, 0:05)")

# =====================================================================
# Slide 9 — Autenticación
# =====================================================================
s = slide()
header(s, 9, "Seguridad", "Autenticación — ¿quién eres?")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    ("Auth delegada a Supabase Auth (no manejamos passwords).", 0),
    ("Cada request trae Authorization: Bearer <JWT>.", 0),
    ("AuthGuard en la API valida antes de cualquier endpoint:", 0),
    ("extrae el token del header", 1),
    ("lo valida contra Supabase (auth.getUser)", 1),
    ("carga el perfil: rol + tenant desde 'profiles'", 1),
    ("inyecta el usuario en el request", 1),
    ("Sin token válido o sin perfil → 401.", 0, RED),
    ("Caché de perfiles en memoria (TTL 5 min).", 0),
], size=14.5, gap=6)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Flujo del token",
     ["1 · Login en Supabase → JWT firmado",
      "2 · Front guarda sesión en cookie",
      "3 · authFetch envía Bearer a la API",
      "4 · AuthGuard valida y resuelve perfil",
      "5 · Endpoint se ejecuta con el usuario",
      "", "No inventamos criptografía propia."],
     accent=BRAND)
notes(s, "Autenticación = saber quién eres. No inventamos sistema de passwords: lo delega Supabase. "
         "El usuario se loguea, recibe un JWT firmado que viaja en Authorization en cada llamada. "
         "El AuthGuard es un portero antes de cada endpoint: saca el token, lo valida contra Supabase, "
         "busca el perfil (rol + tenant). Si falla → 401. Cacheamos perfiles 5 min para no golpear la DB. (1:15)")

# =====================================================================
# Slide 10 — Autorización por roles (tabla)
# =====================================================================
s = slide()
header(s, 10, "Seguridad", "Autorización — ¿qué puedes hacer?")
bullets(s, Inches(0.62), Inches(1.95), Inches(12), Inches(1.0), [
    ("3 roles: admin · jefe_taller · mechanic. RolesGuard + @Roles(...) por endpoint.", 0),
    ("Modelo de mínimo privilegio — reglas en el código, no en la UI.", 0),
], size=15, gap=6)
# tabla
tx, ty = Inches(0.62), Inches(3.15)
rows = [
    ("Acción sobre un ticket", "Roles permitidos", True),
    ("Listar / ver", "admin · jefe_taller · mechanic", False),
    ("Asignar / reasignar", "admin · jefe_taller", False),
    ("Iniciar / finalizar", "mechanic (el asignado)", False),
    ("Validar / cerrar", "admin", False),
]
rh = Inches(0.62); cw1 = Inches(5.4); cw2 = Inches(6.7)
for i, (a, b, hd) in enumerate(rows):
    yy = ty + rh * i
    fill = BRAND_DK if hd else (SURFACE if i % 2 else SURFACE_2)
    rect(s, tx, yy, cw1, rh, fill=fill, line=BORDER, line_w=0.75)
    rect(s, tx + cw1, yy, cw2, rh, fill=fill, line=BORDER, line_w=0.75)
    col = FG if hd else FG
    text(s, tx + Inches(0.2), yy, cw1 - Inches(0.3), rh, a, size=14,
         color=col, bold=hd, anchor=MSO_ANCHOR.MIDDLE)
    text(s, tx + cw1 + Inches(0.2), yy, cw2 - Inches(0.3), rh, b, size=14,
         color=(FG if hd else BRAND_LT), bold=hd, anchor=MSO_ANCHOR.MIDDLE, font=F_MONO if not hd else F_SANS)
notes(s, "Autenticado no es autorizado. Tres roles: admin, jefe de taller, mecánico. Cada endpoint "
         "lleva @Roles que el RolesGuard hace cumplir. Mínimo privilegio: en un ticket, cualquiera lo "
         "ve, pero sólo jefe/admin asigna; sólo el mecánico asignado inicia y finaliza; sólo admin "
         "valida y cierra. Un mecánico no se auto-asigna ni se cierra sus tickets. Reglas en el código, "
         "no en esconder botones. (1:15)")

# =====================================================================
# Slide 11 — Multi-tenant
# =====================================================================
s = slide()
header(s, 11, "Seguridad", "Aislamiento multi-tenant")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.5), [
    ("Cada usuario pertenece a un tenant (taller).", 0),
    ("TenantService resuelve el tenant_id desde el TOKEN…", 0),
    ("…nunca desde lo que mande el cliente.", 0, BRAND_LT),
    ("Toda consulta filtra por tenant_id.", 0),
    ("Un taller jamás ve datos de otro.", 0, GREEN),
], size=16, gap=10)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Ataque que cierra",
     ["Cambiar un ID en la URL para ver datos de otro taller.",
      "",
      "Aunque lo intenten, la API siempre filtra por el tenant de la sesión autenticada.",
      "",
      "El tenant no es un parámetro: se deriva del token."],
     accent=RED, title_color=RED)
notes(s, "Tercer pilar: aislamiento entre clientes. Todo lleva tenant_id. La clave es de dónde sale: "
         "TenantService lo lee del token autenticado, NUNCA de un parámetro del cliente. Eso cierra el "
         "ataque clásico de cambiar un ID en la URL para ver datos de otro taller. La API siempre filtra "
         "por el tenant de tu sesión. (0:50)")

# =====================================================================
# Slide 12 — Evidencias / archivos
# =====================================================================
s = slide()
header(s, 12, "Seguridad", "Evidencias / archivos — signed URLs")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    ("El archivo nunca pasa por nuestra API.", 0, BRAND_LT),
    ("1 · Cliente pide URL → API valida rol, tamaño y tipo", 0),
    ("genera URL firmada (TTL 60 s)", 1),
    ("2 · Cliente sube directo a Supabase Storage", 0),
    ("3 · Cliente confirma → API verifica que existe y registra", 0),
    ("Descarga con URL firmada temporal (TTL 5 min).", 0),
], size=15, gap=7)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Controles",
     ["MIME whitelist: jpg · png · webp",
      "Tamaño máximo: 5 MB",
      "Ruta scopeada: tenant_id/ticket_id/uuid.ext",
      "Descarga firmada y temporal (5 min)",
      "Service-role key SÓLO en el servidor",
      "", "Acceso por rol + asignación al ticket"],
     accent=AMBER, title_color=AMBER)
notes(s, "Las evidencias son fotos de los mecánicos; los archivos siempre son delicados. Patrón de "
         "URLs firmadas: el archivo NO pasa por la API. El cliente pide permiso, la API valida rol, "
         "tipo (sólo imágenes) y tamaño (máx 5MB) y entrega una URL firmada de 60s. Sube directo a "
         "Supabase, confirma, y la API verifica que exista antes de registrar. Descarga: URLs firmadas "
         "que expiran en 5 min, no enlaces públicos. La llave de admin de Supabase vive sólo en el "
         "servidor, nunca en el navegador. (1:15)")

# =====================================================================
# Slide 13 — Secretos, red y RLS
# =====================================================================
s = slide()
header(s, 13, "Seguridad", "Secretos, red y RLS")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    (".env en .gitignore; sólo se versiona .env.example sin valores.", 0),
    ("Service-role key y DB password sólo en el entorno del servidor.", 0),
    ("RLS habilitado en Supabase → segunda capa a nivel de base.", 0),
    ("HTTPS extremo a extremo · SSL obligatorio hacia la DB.", 0),
], size=15, gap=9)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Defensa en profundidad",
     ["Para que un dato salga indebidamente tendría que fallar TODO:",
      "1 · validación de input",
      "2 · guard de auth (JWT)",
      "3 · guard de roles",
      "4 · filtro de tenant",
      "5 · RLS en la base",
      "", "Varias rejas, no una sola."],
     accent=BRAND_LT, title_color=BRAND_LT)
notes(s, "Cierre técnico: los secretos nunca van al repo; sólo subimos una plantilla vacía. Las llaves "
         "viven en Railway/Supabase. Encima, Supabase tiene RLS activado = segunda reja en la base. "
         "Todo es HTTPS. La idea: defensa en profundidad — para filtrar un dato tendría que fallar la "
         "validación Y auth Y roles Y tenant Y RLS. (0:45)")

# =====================================================================
# Slide 14 — Proceso / Gitflow
# =====================================================================
s = slide()
header(s, 14, "Seguridad", "Proceso — Gitflow & PRs")
bullets(s, Inches(0.62), Inches(2.0), Inches(12), Inches(4.0), [
    ("main protegida: nada se mergea sin Pull Request revisado.", 0),
    ("Una rama = un ticket de Linear = un PR · squash merge.", 0),
    ("Conventional Commits + Conventional Branches con ID de Linear.", 0),
    ("Trazabilidad automática: la rama mueve el estado del ticket.", 0),
    ("La seguridad también es proceso: reduce el riesgo de código sin revisar en producción.", 0, BRAND_LT),
], size=16, gap=11)
notes(s, "La seguridad también es proceso. main está protegida: nadie sube directo, todo entra por PR "
         "con revisión. Cada cambio se ata a un ticket de Linear → trazabilidad de quién cambió qué y "
         "por qué. Reduce el riesgo de meter código malo o sin revisar a producción. (0:30)")

# =====================================================================
# Slide 15 — Cierre
# =====================================================================
s = slide()
rect(s, 0, 0, Inches(0.22), EMU_H, fill=BRAND)
text(s, Inches(0.95), Inches(1.0), Inches(11), Inches(0.8), "En resumen", size=34, color=FG, bold=True)
rect(s, Inches(0.98), Inches(1.85), Inches(1.2), Pt(4), fill=BRAND)
card(s, Inches(0.95), Inches(2.25), Inches(5.6), Inches(3.4), "Infraestructura",
     ["3 capas desacopladas:",
      "Next.js (Vercel) · NestJS (Railway) · Supabase",
      "Cada una desplegada y escalada aparte.",
      "El front nunca toca la base: todo por la API."],
     accent=BLUE)
card(s, Inches(6.75), Inches(2.25), Inches(5.85), Inches(3.4), "Seguridad en capas",
     ["validación → JWT → roles → multi-tenant",
      "→ signed URLs → RLS",
      "Sin secretos en el repo · HTTPS · mínimo privilegio",
      "Defensa en profundidad de punta a punta."],
     accent=BRAND_LT, title_color=BRAND_LT)
footer(s, 15)
notes(s, "Resumen: tres capas independientes, mantenibles y escalables. Seguridad en capas: validamos "
         "todo lo que entra, verificamos quién eres (JWT), qué puedes hacer (roles), aislamos por tenant, "
         "archivos con URLs firmadas temporales, y RLS como red final. Sin secretos en el código, todo "
         "cifrado. Abro preguntas. (0:30)")

out = r"C:\Users\danie\Desktop\egenya\trackt-2\trackt-app\documentacion\presentacion-infra-seguridad.pptx"
prs.save(out)
print("OK ->", out, "slides:", len(prs.slides._sldIdLst))
