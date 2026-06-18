# -*- coding: utf-8 -*-
"""Merge: antepone la intro de infra/seguridad de Ramón (12 slides) a la
presentación grupal, elimina las slides 6/17/18/19 y actualiza la agenda.
Las 12 slides se dibujan nativas sobre el canvas del equipo (10x5.625") aplicando
un factor de escala uniforme; reusa la lógica/paleta de build_pptx.py."""

import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches as _Inches, Pt as _Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = r"C:\Users\danie\Desktop\egenya\trackt-2\trackt-app\documentacion\presentacion-trackt.pptx"
OUT = r"C:\Users\danie\Desktop\egenya\trackt-2\trackt-app\documentacion\presentacion-trackt-final.pptx"

prs = Presentation(SRC)
N_TEAM = len(prs.slides._sldIdLst)            # 21
SCALE = prs.slide_width / _Inches(13.333)     # ~0.75 (canvas equipo / diseño 13.333")

# Inches/Pt escalados → todas las coords y fuentes encajan en el canvas del equipo
def Inches(v): return Emu(int(_Inches(v) * SCALE))
def Pt(v):     return _Pt(v * SCALE)

# ---------- Paleta Trackt (dark) ----------
BG        = RGBColor(0x0B, 0x0B, 0x0D)
SURFACE   = RGBColor(0x15, 0x15, 0x18)
SURFACE_2 = RGBColor(0x1D, 0x1D, 0x21)
BORDER    = RGBColor(0x2A, 0x2A, 0x30)
FG        = RGBColor(0xF5, 0xF5, 0xF5)
MUTED     = RGBColor(0xA0, 0xA0, 0xA8)
FAINT     = RGBColor(0x70, 0x70, 0x78)
BRAND     = RGBColor(0x61, 0x52, 0xE8)
BRAND_LT  = RGBColor(0x9B, 0x8F, 0xFF)
BRAND_DK  = RGBColor(0x3D, 0x31, 0xB0)
BLUE   = RGBColor(0x60, 0xA5, 0xFA)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
GREEN  = RGBColor(0x4A, 0xDE, 0x80)
RED    = RGBColor(0xF8, 0x71, 0x71)
EMER   = RGBColor(0x34, 0xD3, 0x99)
F_SANS = "Arial"
F_MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)    # ≈ canvas del equipo
BLANK = prs.slide_layouts[0]                  # layout DEFAULT (sin placeholders)

# ---------- Helpers (portados de build_pptx.py, con Inches/Pt escalados) ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def text(s, x, y, w, h, runs, size=18, color=FG, bold=False, font=F_SANS,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    if spacing: p.line_spacing = spacing
    if isinstance(runs, str): runs = [(runs, {})]
    for t, st in runs:
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(st.get("size", size)); f.bold = st.get("bold", bold)
        f.italic = st.get("italic", italic); f.name = st.get("font", font)
        f.color.rgb = st.get("color", color)
    return tb

def bullets(s, x, y, w, h, items, size=16, gap=8, marker_color=BRAND):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for it in items:
        txt = it[0]; lvl = it[1]; col = it[2] if len(it) > 2 else FG
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.line_spacing = 1.05
        if lvl == 0:
            mk = p.add_run(); mk.text = "▸  "
            mk.font.name = F_SANS; mk.font.size = Pt(size); mk.font.bold = True
            mk.font.color.rgb = marker_color
            r = p.add_run(); r.text = txt
            r.font.name = F_SANS; r.font.size = Pt(size); r.font.color.rgb = col
        else:
            p.level = 1
            mk = p.add_run(); mk.text = "      –   "
            mk.font.name = F_SANS; mk.font.size = Pt(size-1); mk.font.color.rgb = FAINT
            r = p.add_run(); r.text = txt
            r.font.name = F_SANS; r.font.size = Pt(size-1); r.font.color.rgb = MUTED
    return tb

def footer(s, idx):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(7.06), Inches(0.13), Inches(0.13))
    d.fill.solid(); d.fill.fore_color.rgb = BRAND; d.line.fill.background(); d.shadow.inherit = False
    text(s, Inches(0.82), Inches(6.98), Inches(7), Inches(0.3),
         "Trackt · Infraestructura & Seguridad", size=10, color=FAINT)
    text(s, Inches(11.2), Inches(6.98), Inches(1.5), Inches(0.3),
         f"{idx:02d}", size=10, color=FAINT, align=PP_ALIGN.RIGHT)

def header(s, idx, kicker, title, kicker_color=BRAND_LT):
    rect(s, 0, 0, Inches(0.16), EMU_H, fill=BRAND)
    text(s, Inches(0.6), Inches(0.5), Inches(11), Inches(0.3),
         kicker.upper(), size=12.5, color=kicker_color, bold=True)
    text(s, Inches(0.6), Inches(0.82), Inches(12), Inches(0.8),
         title, size=30, color=FG, bold=True, font=F_SANS)
    rect(s, Inches(0.62), Inches(1.62), Inches(0.9), Pt(3), fill=BRAND)
    footer(s, idx)

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def card(s, x, y, w, h, title, lines, accent=BRAND, title_color=FG):
    rect(s, x, y, w, h, fill=SURFACE, line=BORDER, line_w=1.0, rounded=True, radius=0.06)
    rect(s, x, y, w, Pt(4), fill=accent)
    text(s, x + Inches(0.22), y + Inches(0.18), w - Inches(0.4), Inches(0.4),
         title, size=15, color=title_color, bold=True)
    tb = s.shapes.add_textbox(x + Inches(0.22), y + Inches(0.62), w - Inches(0.4), h - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5); p.line_spacing = 1.04
        r = p.add_run(); r.text = ln
        r.font.name = F_SANS; r.font.size = Pt(12.5); r.font.color.rgb = MUTED

def divider(idx, num, title, sub):
    s = slide()
    rect(s, 0, 0, Inches(0.22), EMU_H, fill=BRAND)
    text(s, Inches(0.95), Inches(2.5), Inches(4), Inches(1.2), num, size=90, color=SURFACE_2, bold=True)
    text(s, Inches(0.98), Inches(3.75), Inches(11), Inches(0.9), title, size=44, color=FG, bold=True)
    rect(s, Inches(1.0), Inches(4.75), Inches(1.3), Pt(4), fill=BRAND)
    text(s, Inches(1.0), Inches(4.95), Inches(10.5), Inches(0.6), sub, size=16, color=MUTED)
    footer(s, idx)
    return s

# =====================================================================
# 12 slides de Ramón (intro infra/seguridad) — footers 1..12
# =====================================================================
# 1 — Divider Infraestructura
divider(1, "01", "Infraestructura", "Tres capas independientes: cada una se despliega, escala y falla por separado.")

# 2 — Arquitectura general
s = slide(); header(s, 2, "Infraestructura", "Vista general de la arquitectura")
cy = Inches(2.1); cw = Inches(3.7); ch = Inches(3.05)
card(s, Inches(0.6), cy, cw, ch, "Frontend  ·  SSR",
     ["Next.js 16 (App Router)", "React 19", "TanStack Query", "shadcn/ui + Tailwind 4", "", "Deploy: Vercel"], accent=BLUE)
card(s, Inches(4.82), cy, cw, ch, "Backend  ·  API REST",
     ["NestJS 11 + TypeScript", "Prisma 6 (ORM)", "Auth JWT (Bearer)", "Validación global de DTOs", "", "Deploy: Railway"], accent=BRAND)
card(s, Inches(9.04), cy, cw, ch, "Datos  ·  Servicios",
     ["Supabase", "· PostgreSQL", "· Auth (JWT)", "· Storage (evidencias)", "", "Región: East US"], accent=EMER)
for ax in [Inches(4.42), Inches(8.64)]:
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, cy + ch/2 - Inches(0.16), Inches(0.34), Inches(0.32))
    a.fill.solid(); a.fill.fore_color.rgb = BRAND; a.line.fill.background(); a.shadow.inherit=False
bullets(s, Inches(0.62), Inches(5.55), Inches(12), Inches(1.2), [
    ("El frontend nunca habla directo con la base: todo pasa por la API.", 0),
    ("Comunicación HTTPS + token JWT en cada request → un único punto de entrada controlado.", 0),
], size=15)

# 3 — Frontend
s = slide(); header(s, 3, "Infraestructura", "Frontend — Next.js 16 + Vercel")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("Next.js 16 (App Router) + React 19, rendering en servidor.", 0),
    ("Protección de rutas en el servidor:", 0),
    ("middleware valida la sesión (supabase.auth.getUser) en cada request", 1),
    ("redirige a /login si no hay sesión, antes de renderizar", 1),
    ("Sesión en cookies vía @supabase/ssr (no localStorage → menos XSS).", 0),
    ("TanStack Query para estado de servidor y caché.", 0),
    ("shadcn + Tailwind 4 · formularios validados con zod.", 0),
    ("authFetch adjunta el token como Bearer (auto-refresh + retry en 401).", 0),
    ("El build de producción falla si faltan las env vars de Supabase.", 0),
], size=15, gap=7)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Por qué importa",
     ["Sin lógica de negocio ni llaves secretas en el cliente.",
      "Protección server-side: no depende de 'esconder un botón'.",
      "Cookies > localStorage para la sesión.",
      "Token refrescado de forma transparente.", "",
      "producto/tract-front  (paquete trackt-front)"], accent=BLUE)

# 4 — Backend
s = slide(); header(s, 4, "Infraestructura", "Backend / API — NestJS 11")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("NestJS 11 + TypeScript: arquitectura modular (un módulo por dominio).", 0),
    ("Prisma 6 como ORM → PostgreSQL (queries tipadas, anti-inyección).", 0),
    ("Validación global con ValidationPipe:", 0),
    ("whitelist → descarta campos no declarados en el DTO", 1),
    ("transform → convierte y valida tipos antes de la lógica", 1),
    ("Ningún dato entra a la API sin pasar por un filtro.", 0),
    ("Deploy en Railway (Node).", 0),
], size=15, gap=8)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Módulos de dominio",
     ["auth · tenant", "equipos · usuarios", "ordenes · tickets", "evidencias",
      "inventario", "notificaciones", "", "9 módulos aislados → reglas de seguridad por módulo"], accent=BRAND)

# 5 — Base de datos
s = slide(); header(s, 5, "Infraestructura", "Base de datos — Supabase PostgreSQL")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.1), Inches(4.6), [
    ("PostgreSQL administrada por Supabase (región East US).", 0),
    ("Conexión vía pooler (modo transacción, 6543) en runtime.", 0),
    ("Conexión directa (5432) sólo para migraciones.", 0),
    ("sslmode=require → tráfico a la DB siempre cifrado.", 0),
    ("tenant_id en TODAS las tablas de negocio → multi-tenant desde el diseño.", 0, BRAND_LT),
], size=15, gap=9)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Modelo de datos",
     ["Tenant · Equipo", "OrdenTrabajo · Ticket", "Evidencia", "Repuesto · InventarioStock",
      "ReservaRepuesto", "MovimientoInventario", "Notificacion", "", "Cubre todo el flujo del taller"], accent=EMER)

# 6 — Divider Seguridad
divider(6, "02", "Seguridad", "Defensa en profundidad: varias capas, no una sola.")

# 7 — Autenticación
s = slide(); header(s, 7, "Seguridad", "Autenticación — ¿quién eres?")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    ("Auth delegada a Supabase Auth (no manejamos passwords).", 0),
    ("Cada request trae Authorization: Bearer <JWT>.", 0),
    ("AuthGuard en la API valida antes de cualquier endpoint:", 0),
    ("extrae el token del header", 1),
    ("lo valida contra Supabase (auth.getUser)", 1),
    ("carga el perfil: rol + tenant desde 'profiles'", 1),
    ("inyecta el usuario en el request", 1),
    ("Sin token válido o sin perfil → 401.", 0, RED),
    ("Caché de perfiles en memoria (TTL 5 min).", 0),
], size=14.5, gap=6)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Flujo del token",
     ["1 · Login en Supabase → JWT firmado", "2 · Front guarda sesión en cookie",
      "3 · authFetch envía Bearer a la API", "4 · AuthGuard valida y resuelve perfil",
      "5 · Endpoint se ejecuta con el usuario", "", "No inventamos criptografía propia."], accent=BRAND)

# 8 — Roles (tabla)
s = slide(); header(s, 8, "Seguridad", "Autorización — ¿qué puedes hacer?")
bullets(s, Inches(0.62), Inches(1.95), Inches(12), Inches(1.0), [
    ("3 roles: admin · jefe_taller · mechanic. RolesGuard + @Roles(...) por endpoint.", 0),
    ("Modelo de mínimo privilegio — reglas en el código, no en la UI.", 0),
], size=15, gap=6)
tx, ty = Inches(0.62), Inches(3.15)
rows = [
    ("Acción sobre un ticket", "Roles permitidos", True),
    ("Listar / ver", "admin · jefe_taller · mechanic", False),
    ("Asignar / reasignar", "admin · jefe_taller", False),
    ("Iniciar / finalizar", "mechanic (el asignado)", False),
    ("Validar / cerrar", "admin", False),
]
rh = Inches(0.62); cw1 = Inches(5.4); cw2 = Inches(6.7)
for i, (a, b, hd) in enumerate(rows):
    yy = ty + rh * i
    fill = BRAND_DK if hd else (SURFACE if i % 2 else SURFACE_2)
    rect(s, tx, yy, cw1, rh, fill=fill, line=BORDER, line_w=0.75)
    rect(s, tx + cw1, yy, cw2, rh, fill=fill, line=BORDER, line_w=0.75)
    text(s, tx + Inches(0.2), yy, cw1 - Inches(0.3), rh, a, size=14, color=FG, bold=hd, anchor=MSO_ANCHOR.MIDDLE)
    text(s, tx + cw1 + Inches(0.2), yy, cw2 - Inches(0.3), rh, b, size=14,
         color=(FG if hd else BRAND_LT), bold=hd, anchor=MSO_ANCHOR.MIDDLE, font=F_MONO if not hd else F_SANS)

# 9 — Multi-tenant
s = slide(); header(s, 9, "Seguridad", "Aislamiento multi-tenant")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.5), [
    ("Cada usuario pertenece a un tenant (taller).", 0),
    ("TenantService resuelve el tenant_id desde el TOKEN…", 0),
    ("…nunca desde lo que mande el cliente.", 0, BRAND_LT),
    ("Toda consulta filtra por tenant_id.", 0),
    ("Un taller jamás ve datos de otro.", 0, GREEN),
], size=16, gap=10)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Ataque que cierra",
     ["Cambiar un ID en la URL para ver datos de otro taller.", "",
      "Aunque lo intenten, la API siempre filtra por el tenant de la sesión autenticada.", "",
      "El tenant no es un parámetro: se deriva del token."], accent=RED, title_color=RED)

# 10 — Evidencias
s = slide(); header(s, 10, "Seguridad", "Evidencias / archivos — signed URLs")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    ("El archivo nunca pasa por nuestra API.", 0, BRAND_LT),
    ("1 · Cliente pide URL → API valida rol, tamaño y tipo", 0),
    ("genera URL firmada (TTL 60 s)", 1),
    ("2 · Cliente sube directo a Supabase Storage", 0),
    ("3 · Cliente confirma → API verifica que existe y registra", 0),
    ("Descarga con URL firmada temporal (TTL 5 min).", 0),
], size=15, gap=7)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Controles",
     ["MIME whitelist: jpg · png · webp", "Tamaño máximo: 5 MB",
      "Ruta scopeada: tenant_id/ticket_id/uuid.ext", "Descarga firmada y temporal (5 min)",
      "Service-role key SÓLO en el servidor", "", "Acceso por rol + asignación al ticket"],
     accent=AMBER, title_color=AMBER)

# 11 — Secretos / RLS
s = slide(); header(s, 11, "Seguridad", "Secretos, red y RLS")
bullets(s, Inches(0.62), Inches(2.0), Inches(7.0), Inches(4.6), [
    (".env en .gitignore; sólo se versiona .env.example sin valores.", 0),
    ("Service-role key y DB password sólo en el entorno del servidor.", 0),
    ("RLS habilitado en Supabase → segunda capa a nivel de base.", 0),
    ("HTTPS extremo a extremo · SSL obligatorio hacia la DB.", 0),
], size=15, gap=9)
card(s, Inches(8.05), Inches(2.0), Inches(4.7), Inches(4.5), "Defensa en profundidad",
     ["Para que un dato salga indebidamente tendría que fallar TODO:",
      "1 · validación de input", "2 · guard de auth (JWT)", "3 · guard de roles",
      "4 · filtro de tenant", "5 · RLS en la base", "", "Varias rejas, no una sola."],
     accent=BRAND_LT, title_color=BRAND_LT)

# 12 — Gitflow
s = slide(); header(s, 12, "Seguridad", "Proceso — Gitflow & PRs")
bullets(s, Inches(0.62), Inches(2.0), Inches(12), Inches(4.0), [
    ("main protegida: nada se mergea sin Pull Request revisado.", 0),
    ("Una rama = un ticket de Linear = un PR · squash merge.", 0),
    ("Conventional Commits + Conventional Branches con ID de Linear.", 0),
    ("Trazabilidad automática: la rama mueve el estado del ticket.", 0),
    ("La seguridad también es proceso: reduce el riesgo de código sin revisar en producción.", 0, BRAND_LT),
], size=16, gap=11)

# =====================================================================
# Actualizar la AGENDA del equipo (antes de reordenar)
# =====================================================================
# Nuevo orden de agenda: 01 Contexto · 02 Ramón Infra/Seg · 03 Jaime API · 04 Rosio Front · 05 Cierre
REPL = {
    "Rosio — Frontend": "Ramón — Infraestructura & Seguridad",
    "Jaime — Backend": "Jaime — Backend / API",
    "Ramón — Arq. y Datos": "Rosio — Frontend",
    "Arquitectura Next.js · Auth/App · Nav por rol · Design System · Pantallas":
        "Arquitectura · Auth · Roles · Multi-tenant · Evidencias · RLS",
    "Arquitectura 3 capas · DB · RLS · Pruebas · GitHub · Deploy":
        "Arquitectura Next.js · Auth/App · Nav por rol · Design System · Pantallas",
}

def set_text_keep(shape, new):
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = new
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.text = new

agenda = None
for sl in prs.slides:
    txts = [sh.text_frame.text for sh in sl.shapes if sh.has_text_frame]
    if any("Defensa de 30 minutos" in t for t in txts):
        agenda = sl; break
assert agenda is not None, "No se encontró el slide de agenda"
n_repl = 0
for sh in agenda.shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text.strip()
    if t in REPL:
        set_text_keep(sh, REPL[t]); n_repl += 1
print("agenda: reemplazos aplicados =", n_repl, "(esperado 5)")

# =====================================================================
# Conclusiones: quitar "El mayor riesgo..." y los pendientes ①②③④,
# repurpose la card derecha como "Futuras mejoras" (estilo nativo del equipo)
# =====================================================================
def first_run_props(shape):
    r = shape.text_frame.paragraphs[0].runs[0]
    f = r.font
    color = f.color.rgb if (f.color is not None and f.color.type is not None) else None
    return f.name, f.size, f.bold, color

def set_single(shape, txt):
    tf = shape.text_frame
    for p in tf.paragraphs[1:]:          # borrar párrafos extra (texto multilínea)
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    p.runs[0].text = txt
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)

def set_multi(shape, lines):
    name, size, bold, color = first_run_props(shape)
    tf = shape.text_frame
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    def style(run):
        run.font.name = name
        if size: run.font.size = size
        run.font.bold = bold
        if color is not None: run.font.color.rgb = color
    def add_line(par, txt):
        par.space_after = _Pt(7); par.line_spacing = 1.05
        rr = par.add_run(); rr.text = txt; style(rr)
    add_line(p0, lines[0])
    for ln in lines[1:]:
        add_line(tf.add_paragraph(), ln)

concl = None
for sl in prs.slides:
    blob = " ".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
    if "CONCLUSIONES" in blob and "El mayor riesgo" in blob:
        concl = sl; break
assert concl is not None, "No se encontró slide de Conclusiones"
edits = 0
for sh in concl.shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text.strip()
    if t.startswith("El mayor riesgo"):
        set_single(sh, "Futuras mejoras"); edits += 1
    elif t.startswith("① Backup SQL"):
        set_multi(sh, [
            "① Agente de IA — gestión de tickets y operación vía chat",
            "② Notificaciones email + WhatsApp (Kapso)",
            "③ QR para identificar equipos / máquinas en terreno",
        ]); edits += 1
print("conclusiones: ediciones aplicadas =", edits, "(esperado 2)")

# =====================================================================
# Reordenar (12 nuevas al frente) y eliminar slides 6,17,18,19 del equipo
# =====================================================================
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)                       # 33
team, new = ids[:N_TEAM], ids[N_TEAM:]     # 21 (equipo, orden original) + 12 (Ramón)
def T(*ones): return [team[o - 1] for o in ones]   # 1-based del equipo
order = (
    T(1, 2, 3, 4, 5)                # intro: portada · agenda · qué cambió · MVP · tres roles
    + new                           # infra & seguridad (Ramón, 12 slides)
    + T(14, 15, 16)                 # API / backend (Jaime)
    + T(7, 8, 9, 10, 11, 12, 13)    # frontend (Rosio)
    + T(20, 21)                     # cierre: cumplimiento + conclusiones
)                                   # elimina 6,17,18,19 (no incluidas) → 29
for e in ids: sldIdLst.remove(e)
for e in order: sldIdLst.append(e)

cands = [OUT, OUT.replace(".pptx", "-v2.pptx"), OUT.replace(".pptx", "-v3.pptx")]
saved = None
for c in cands:
    try:
        prs.save(c); saved = c; break
    except PermissionError:
        print("  bloqueado (abierto en PowerPoint?):", c)
print("OK ->", saved)
print("slides finales:", len(prs.slides._sldIdLst), "(esperado 29)")
